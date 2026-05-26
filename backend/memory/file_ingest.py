from __future__ import annotations

from pathlib import Path

# 纯文本格式，直接 read_text 即可
_PLAINTEXT_EXTS = {
    ".md", ".markdown", ".txt", ".py", ".js", ".ts", ".jsx", ".tsx",
    ".html", ".htm", ".css", ".json", ".yaml", ".yml", ".toml",
    ".xml", ".rst", ".org", ".tex", ".log", ".ini", ".cfg", ".env",
}


def extract_text(file_path: str) -> str:
    """
    从各种格式文件中提取纯文本。
    支持：PDF / Word(docx) / Excel(xlsx/csv) / PowerPoint(pptx)
         / EPUB / RTF / 纯文本族（md/txt/py/js/...）
    未知格式：尝试 UTF-8 读取，实在不行才报错。
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)

    if ext in _PLAINTEXT_EXTS:
        return path.read_text(encoding="utf-8", errors="replace")

    if ext == ".docx":
        return _extract_docx(file_path)

    if ext in (".xlsx", ".xls"):
        return _extract_excel(file_path)

    if ext == ".csv":
        return _extract_csv(file_path)

    if ext == ".pptx":
        return _extract_pptx(file_path)

    if ext == ".epub":
        return _extract_epub(file_path)

    if ext == ".rtf":
        return _extract_rtf(file_path)

    # 最后兜底：尝试当纯文本读
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        if text.strip():
            return text
    except Exception:
        pass

    raise ValueError(
        f"不支持的文件格式: {ext}。"
        f"支持的格式：PDF, Word(docx), Excel(xlsx/csv), PowerPoint(pptx), "
        f"EPUB, RTF, 以及纯文本类（md/txt/py/js/html/json/yaml等）"
    )


# ── PDF ──────────────────────────────────────────────────────────────────────

def _extract_pdf(file_path: str) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages)
    except ImportError:
        pass
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            pages = [page.extract_text() or "" for page in pdf.pages]
        return "\n\n".join(pages)
    except ImportError:
        raise ImportError("PDF 解析需要安装 PyPDF2 或 pdfplumber")


# ── Word (.docx) ─────────────────────────────────────────────────────────────

def _extract_docx(file_path: str) -> str:
    try:
        import docx
        doc = docx.Document(file_path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # 也提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n\n".join(parts)
    except ImportError:
        raise ImportError("解析 .docx 需要安装 python-docx: pip install python-docx")


# ── Excel (.xlsx / .xls) ─────────────────────────────────────────────────────

def _extract_excel(file_path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"## 工作表: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(v) for v in row if v is not None)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except ImportError:
        pass
    # xlrd 兜底（旧 .xls）
    try:
        import xlrd
        wb = xlrd.open_workbook(file_path)
        parts = []
        for sheet in wb.sheets():
            parts.append(f"## 工作表: {sheet.name}")
            for i in range(sheet.nrows):
                row_text = " | ".join(str(v) for v in sheet.row_values(i) if v != "")
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except ImportError:
        raise ImportError("解析 Excel 需要安装 openpyxl: pip install openpyxl")


# ── CSV ──────────────────────────────────────────────────────────────────────

def _extract_csv(file_path: str) -> str:
    import csv
    rows = []
    # 尝试多种编码
    for encoding in ("utf-8", "gbk", "latin-1"):
        try:
            with open(file_path, newline="", encoding=encoding, errors="replace") as f:
                reader = csv.reader(f)
                rows = [" | ".join(row) for row in reader if any(row)]
            break
        except UnicodeDecodeError:
            continue
    return "\n".join(rows)


# ── PowerPoint (.pptx) ───────────────────────────────────────────────────────

def _extract_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"## 第 {i} 页\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except ImportError:
        raise ImportError("解析 .pptx 需要安装 python-pptx: pip install python-pptx")


# ── EPUB ─────────────────────────────────────────────────────────────────────

def _extract_epub(file_path: str) -> str:
    try:
        import ebooklib
        from ebooklib import epub
        from html.parser import HTMLParser

        class _TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.texts = []
            def handle_data(self, data):
                if data.strip():
                    self.texts.append(data.strip())

        book = epub.read_epub(file_path)
        parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            parser = _TextExtractor()
            parser.feed(item.get_content().decode("utf-8", errors="replace"))
            if parser.texts:
                parts.append(" ".join(parser.texts))
        return "\n\n".join(parts)
    except ImportError:
        raise ImportError("解析 .epub 需要安装 ebooklib: pip install ebooklib")


# ── RTF ──────────────────────────────────────────────────────────────────────

def _extract_rtf(file_path: str) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        content = Path(file_path).read_text(encoding="utf-8", errors="replace")
        return rtf_to_text(content)
    except ImportError:
        pass
    # 兜底：简单去除 RTF 控制字符，返回可读文本
    import re
    content = Path(file_path).read_text(encoding="utf-8", errors="replace")
    text = re.sub(r"\\[a-z]+\d* ?", " ", content)
    text = re.sub(r"[{}\\]", "", text)
    return re.sub(r" {2,}", " ", text).strip()
